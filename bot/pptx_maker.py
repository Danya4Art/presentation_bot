from pptx import Presentation
from pptx.util import Inches as In
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
import copy
from stuff.paths import template_path


class PPTXMaker:

    img_top = In(2)
    img_left = In(0.5)
    img_height = In(4.95)
    img_width = In(5.5)
    pgs = [
        'Дата инспекции:',
        'GPS-координаты:',
        'Выводы и замечания:'
    ]

    def __init__(self, template=None, path=None):
        if template is None:
            template = template_path
        self.prs = Presentation(template)
        self.path = path

    @property
    def slides_count(self):
        return len(self.prs.slides)

    def create_slide(self, layout_type=3):
        self.prs.slides.add_slide(self.prs.slide_layouts[layout_type])
        # self.prs.slides[-1].follow_master_background = True
        if layout_type == 3:
            self.prs.slides[-1].shapes[0].top = In(1)
            self.prs.slides[-1].shapes[0].left = In(0.5)
            self.prs.slides[-1].shapes[0].height = In(1)
            self.prs.slides[-1].shapes[0].width = In(9)
            self.prs.slides[-1].shapes[1].top = self.img_top
            self.prs.slides[-1].shapes[1].left = self.img_left + self.img_width
            self.prs.slides[-1].shapes[1].height = self.img_height
            self.prs.slides[-1].shapes[1].width = self.img_width - In(2)
            self.delete_shape(-1, 2)
        self.prs.slides[-1].shapes.add_textbox(8688960, 6503760, width=In(0.1), height=In(0.1))
        self.prs.slides[-1].shapes[-1].text = str(len(self.prs.slides))
        for par in self.pgs:
            self.put_text(-1, 1, par, bold=True)
            self.put_text(-1, 1, '', level=1)

    def put_text(self, sld, shape, text, paragraph=None, size=16, center=False, bold=False, level=0):
        tf = self.prs.slides[sld].shapes[shape].text_frame
        if paragraph is None:
            paragraph = -1
            if tf.paragraphs[0].text:
                tf.add_paragraph()
        tf.paragraphs[paragraph].text = text
        tf.paragraphs[paragraph].font.size = Pt(size)
        tf.paragraphs[paragraph].font.bold = bold
        tf.paragraphs[paragraph].level = level
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        if center:
            tf.paragraphs[paragraph].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    def put_image(self, sld, img):
        images = list(filter(lambda shape: hasattr(shape, 'image'), self.prs.slides[sld].shapes))
        left, top, height, width = self.img_left, self.img_top, self.img_height, self.img_width
        if len(images) == 1:
            images[0].width //= 2
            images[0].height //= 2
            width //= 2
            height //= 2
            left += width
        elif len(images) == 2:
            width //= 2
            height //= 2
            top += height
        elif len(images) == 3:
            width //= 2
            height //= 2
            left += width
            top += height
        elif len(images) > 3:
            raise OverflowError('Too many pictures on one slide')
        self.prs.slides[sld].shapes.add_picture(img, left, top, width=width, height=height)

    def delete_shape(self, sld, shape):
        elem = self.prs.slides[sld].shapes[shape]._element
        elem.getparent().remove(elem)

    def delete_slide(self, sld):
        xml_slides = self.prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[sld])

    def save(self, spec_path=None):
        if spec_path is not None:
            self.prs.save(spec_path)
        elif self.path is not None:
            self.prs.save(self.path)
        else:
            print('ПРЕЗЕНТАЦИЯ НЕ БЫЛА СОХРАНЕНА, ДАННЫЕ УТЕРЯНЫ')
            raise IOError

    def get_content(self, sld):
        content = []
        head = 0 if sld else 2
        images = list(filter(lambda shape: hasattr(shape, 'image'), self.prs.slides[sld].shapes))
        content.append(f'Слайд номер {sld + 1} из {self.slides_count}')
        content.append(f'Заголовок:\n{self.prs.slides[sld].shapes[head].text}')
        if sld:
            content.append(f'Текст:\n{self.prs.slides[sld].shapes[1].text}')
            content.append(f'Количество изображений: {len(images)}')
        return content

    def duplicate_slide(self, index=-1):
        # if index == -1:
        #     index = self.slides_count - 1
        # template = self.prs.slides[index]
        # try:
        #     blank_slide_layout = self.prs.slide_layouts[0]
        # except:
        #     blank_slide_layout = self.prs.slide_layouts[len(self.prs.slide_layouts)]
        # copied_slide = self.prs.slides.add_slide(blank_slide_layout)
        # for shp in filter(lambda s: not hasattr(s, 'image'), template.shapes):
        #     el = shp.element
        #     newel = copy.deepcopy(el)
        #     copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        main_slide = self.prs.slides[index]
        self.create_slide()
        self.put_text(-1, 0, main_slide.shapes[0].text, size=18, center=True, bold=True)
        for i, par in enumerate(main_slide.shapes[1].text_frame.paragraphs):
            self.prs.slides[-1].shapes[1].text_frame.paragraphs[i].text = par.text
        # self.prs.slides[-1].shapes[-1].text = str(len(self.prs.slides))
        # for _, value in six.iteritems(template.part.rels):
        #     # Make sure we don't copy a notesSlide relation as that won't exist
        #     if "notesSlide" not in value.reltype:
        #         copied_slide.part.rels.add_relationship(
        #             value.reltype,
        #             value._target,
        #             value.rId
        #         )


# pm = PPTXMaker('template.pptx', '../test.pptx')

